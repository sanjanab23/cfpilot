"""
ppt_generator.py  ── v2.2
Builds a Tata-branded PPTX from a list of slide JSON objects.

UPGRADES:
  V1  Removed accent underline under titles (biggest AI-slide tell)
  V2  Full-width dark header band on every content slide
  V4  Layout rotation (3 alternating layouts — full / panel-right / dark-sidebar)
  V5  Insight chips with coloured left-bar instead of plain ▸ bullets
  V6  Dark mode cover slide with diagonal geometric accents
  V7  Section slide: ghost watermark number behind title
  V8  KPI cards: bottom accent strip + icon circle for depth
  V9  Slide number + progress dots in bottom bar  (via add_bottom_bar)
  V10 Thank You: geometric background shapes
  V11 KPI slide: insight chips rendered below KPI cards (fix)
  V12 So What / bullets slide: insight chips rendered when insights present
  E4  Watermark mode (DRAFT / CONFIDENTIAL diagonal stamp)
  P6  Fallback bullets slide so deck is never broken on build error
  S1  _sanitize_slide_dict applied at generate_pptx level
  NO_FUNNEL  Pyramid/funnel slide removed
"""

import io
import logging
from typing import List, Dict, Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from ppt_utils import (
    SLIDE_W, SLIDE_H,
    TATA_BLUE, TATA_DARK_BLUE, TATA_LIGHT_BLUE, TATA_WHITE,
    TATA_OFF_WHITE, TATA_TEXT_DARK, TATA_TEXT_MID, TATA_TEAL,
    add_logo, add_logo_top_right, add_bottom_bar,
    chart_to_png, timeline_to_png,
    has_enough_data, get_logo_bytes,
)

logger = logging.getLogger(__name__)

# ─── V4: LAYOUT ROTATION ──────────────────────────────────────────────────────
_LAYOUTS = ["full", "panel-right", "dark-sidebar"]
_layout_idx = 0

def _next_layout() -> str:
    global _layout_idx
    l = _LAYOUTS[_layout_idx % len(_LAYOUTS)]
    _layout_idx += 1
    return l

def _reset_layouts() -> None:
    global _layout_idx
    _layout_idx = 0

# ─── CORE PRIMITIVES ──────────────────────────────────────────────────────────

def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def _rect(slide, l, t, w, h, color: RGBColor):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def _txt(slide, l, t, w, h, text: str,
         size: float = 14, bold: bool = False, italic: bool = False,
         color: RGBColor = None, align=PP_ALIGN.LEFT,
         wrap: bool = True, font: str = "Calibri"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.auto_size  = None
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text       = text
    run.font.name  = font
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return tb

def _img(slide, img_bytes: bytes, l, t, w, h) -> None:
    slide.shapes.add_picture(io.BytesIO(img_bytes), l, t, w, h)

# ─── V2: DARK HEADER BAND ─────────────────────────────────────────────────────

HEADER_H   = Inches(1.01)
SEP_LINE_H = Inches(0.012)
TATA_RED   = RGBColor(0xE0, 0x20, 0x3C)

def _header_band(slide, title: str, subtitle: str = "") -> None:
    _rect(slide, Inches(0), Inches(0), SLIDE_W - Inches(2.2), HEADER_H, RGBColor(0x09, 0x41, 0x7A))
    _rect(slide, Inches(0), HEADER_H, SLIDE_W, SEP_LINE_H, RGBColor(0xBB, 0xBB, 0xBB))
    _txt(slide, Inches(0.28), Inches(0.11),
         SLIDE_W - Inches(3.2), Inches(0.55),
         title, size=22, bold=True, color=TATA_WHITE, font="Calibri")
    if subtitle:
        _txt(slide, Inches(0.28), Inches(0.63),
             SLIDE_W - Inches(3.2), Inches(0.33),
             subtitle, size=11, bold=False, color=RGBColor(0xBB, 0xDE, 0xFB),
             font="Calibri")
    _header_logo(slide)


def _header_logo(slide) -> None:
    from ppt_utils import get_logo_bytes
    logo = get_logo_bytes()
    if not logo:
        return
    try:
        lw = Inches(1.96)
        lh = Inches(0.72)
        ll = SLIDE_W - lw - Inches(0.18)
        lt = (HEADER_H - lh) / 2
        slide.shapes.add_picture(io.BytesIO(logo), ll, lt, lw, lh)
    except Exception as e:
        logger.warning(f"_header_logo failed: {e}")

# ─── V5: INSIGHT CHIPS ───────────────────────────────────────────────────────

def _insight_chips(slide, insights: List[str], l, t, w, h) -> None:
    """V5 ── Coloured left-bar chip for each insight bullet."""
    if not insights:
        return
    chip_h   = Inches(0.54)
    chip_gap = Inches(0.10)
    bar_w    = Inches(0.055)
    chip_colors = [TATA_BLUE, TATA_TEAL, RGBColor(0x7B, 0x1F, 0xA2)]
    for i, insight in enumerate(insights[:3]):
        ct = t + i * (chip_h + chip_gap)
        if ct + chip_h > t + h:
            break
        color = chip_colors[i % len(chip_colors)]
        _rect(slide, l, ct, w, chip_h, RGBColor(0xF0, 0xF5, 0xFF))
        _rect(slide, l, ct, bar_w, chip_h, color)
        _txt(slide, l + bar_w + Inches(0.12), ct + Inches(0.08),
             w - bar_w - Inches(0.18), chip_h - Inches(0.12),
             insight, size=10.5, color=TATA_TEXT_DARK, wrap=True)

# ─── E4: WATERMARK ────────────────────────────────────────────────────────────

def _watermark(slide, text: str = "DRAFT") -> None:
    """E4 ── Diagonal semi-transparent stamp across the slide."""
    try:
        tb = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10), Inches(4))
        tf = tb.text_frame
        tf.word_wrap = False
        p   = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text            = text
        run.font.size       = Pt(96)
        run.font.bold       = True
        run.font.color.rgb  = RGBColor(0xE0, 0xE0, 0xE0)
        run.font.name       = "Calibri"
        tb.rotation         = -45
    except Exception as e:
        logger.warning(f"Watermark failed: {e}")

# ─── SLIDE BUILDERS ───────────────────────────────────────────────────────────

def build_cover(prs: Presentation, d: Dict):
    """V6 ── Dark mode cover with diagonal geometric accents."""
    slide = _blank(prs)
    _set_bg(slide, TATA_DARK_BLUE)

    for (l, t, w, h, rot, color) in [
        (Inches(8.5),  Inches(4.2), Inches(6.0), Inches(5.0),  30, TATA_BLUE),
        (Inches(10.5), Inches(-0.5), Inches(4.0), Inches(4.0), 15, RGBColor(0x0A,0x3A,0x80)),
    ]:
        s = slide.shapes.add_shape(1, l, t, w, h)
        s.fill.solid(); s.fill.fore_color.rgb = color
        s.line.fill.background(); s.rotation = rot

    _rect(slide, Inches(0), Inches(0), Inches(0.22), SLIDE_H, RGBColor(0x03,0xA9,0xF4))

    add_logo_top_right(slide, size_inches=1.0)

    _txt(slide, Inches(0.5), Inches(0.28), Inches(8.0), Inches(0.45),
         "TATA CHEMICALS LIMITED", size=10, bold=True,
         color=RGBColor(0x90,0xBC,0xEA), font="Calibri")
    _rect(slide, Inches(0.5), Inches(0.78), Inches(7.0), Inches(0.03),
          RGBColor(0x03,0xA9,0xF4))

    _txt(slide, Inches(0.5), Inches(2.1), Inches(9.5), Inches(2.2),
         d.get("title", "Presentation"),
         size=38, bold=True, color=TATA_WHITE, font="Calibri")

    sub = d.get("subtitle", "")
    if sub:
        _txt(slide, Inches(0.5), Inches(4.45), Inches(9.0), Inches(0.7),
             sub, size=14, color=RGBColor(0xBB,0xDE,0xFB), font="Calibri")

    _rect(slide, Inches(0), SLIDE_H - Inches(0.5), SLIDE_W, Inches(0.5),
          RGBColor(0x0A,0x3A,0x80))
    _rect(slide, Inches(0), SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), TATA_BLUE)
    return slide


def build_section(prs: Presentation, d: Dict, num: int = 1):
    """V7 ── Ghost watermark number behind the section title."""
    slide = _blank(prs)
    _set_bg(slide, TATA_OFF_WHITE)

    bw = Inches(5.3)
    _rect(slide, Inches(0), Inches(0), bw, SLIDE_H, TATA_DARK_BLUE)

    _txt(slide, Inches(-0.5), Inches(0.5), bw + Inches(0.5), Inches(6.0),
         f"0{num}", size=220, bold=True,
         color=RGBColor(0x1A,0x5A,0xA8), font="Calibri")

    _txt(slide, Inches(0.45), Inches(1.2), Inches(1.5), Inches(1.2),
         f"{num:02d}", size=52, bold=True,
         color=RGBColor(0x03,0xA9,0xF4), font="Calibri")
    _txt(slide, Inches(0.45), Inches(2.9), bw - Inches(0.6), Inches(1.8),
         d.get("title", "Section"),
         size=28, bold=True, color=TATA_WHITE, font="Calibri")

    _rect(slide, bw + Inches(0.5), Inches(2.6), Inches(0.05), Inches(2.5), TATA_BLUE)

    sub = d.get("subtitle", "")
    if sub:
        _txt(slide, bw + Inches(0.7), Inches(2.9),
             SLIDE_W - bw - Inches(0.9), Inches(2.0),
             sub, size=14, color=TATA_TEXT_MID, font="Calibri")

    add_bottom_bar(slide)
    return slide


def build_kpi(prs: Presentation, d: Dict,
              slide_num: int = None, total: int = None):
    """
    V8 ── KPI cards with bottom accent strip + icon circle.
    V11 ── Insight chips rendered below cards when insights are present.
    """
    slide    = _blank(prs)
    _set_bg(slide, TATA_OFF_WHITE)
    data     = d.get("data", [])
    insights = d.get("insights", [])

    _header_band(slide, d.get("title", "Key Metrics"))

    card_colors  = [TATA_BLUE, TATA_DARK_BLUE, TATA_TEAL, RGBColor(0x7B,0x1F,0xA2)]
    strip_colors = [RGBColor(0x0D,0x47,0xA1), RGBColor(0x07,0x2F,0x6E),
                    RGBColor(0x00,0x5F,0x6B), RGBColor(0x56,0x12,0x75)]

    n      = min(len(data), 4)
    card_w = (SLIDE_W - Inches(1.0)) / max(n, 1) - Inches(0.2)

    # V11 — shrink card height to leave room for insight chips when present
    has_insights = bool(insights)
    card_h = Inches(2.3) if has_insights else Inches(2.6)
    card_t = Inches(1.25)
    gap    = Inches(0.22)

    for i, item in enumerate(data[:4]):
        cl = Inches(0.5) + i * (card_w + gap)
        _rect(slide, cl, card_t, card_w, card_h, card_colors[i % 4])
        # Bottom strip (V8)
        sh = card_h * 0.22
        _rect(slide, cl, card_t + card_h - sh, card_w, sh, strip_colors[i % 4])
        # Icon circle (V8)
        cs   = Inches(0.32)
        circ = slide.shapes.add_shape(9, cl + Inches(0.18), card_t + Inches(0.18), cs, cs)
        circ.fill.solid(); circ.fill.fore_color.rgb = TATA_WHITE; circ.line.fill.background()
        inn  = slide.shapes.add_shape(9, cl + Inches(0.25), card_t + Inches(0.25),
                                      Inches(0.18), Inches(0.18))
        inn.fill.solid(); inn.fill.fore_color.rgb = card_colors[i % 4]; inn.line.fill.background()
        # Value
        _txt(slide, cl + Inches(0.12), card_t + Inches(0.55),
             card_w - Inches(0.22), Inches(1.0),
             str(item.get("value", "—")),
             size=34, bold=True, color=TATA_WHITE, align=PP_ALIGN.CENTER)
        # Label
        _txt(slide, cl + Inches(0.1), card_t + Inches(1.65),
             card_w - Inches(0.2), Inches(0.55),
             str(item.get("label", "")),
             size=10.5, color=RGBColor(0xBB,0xDE,0xFB), align=PP_ALIGN.CENTER)

    # V11 ── Insight chips below KPI cards
    if has_insights:
        chips_top = card_t + card_h + Inches(0.18)
        chips_h   = SLIDE_H - chips_top - Inches(0.18)
        _insight_chips(slide, insights,
                       Inches(0.5), chips_top,
                       SLIDE_W - Inches(1.0), chips_h)

    add_bottom_bar(slide, slide_num, total)
    return slide


def build_chart(prs: Presentation, d: Dict,
                slide_num: int = None, total: int = None):
    """V4 ── Rotates across 3 layouts per call."""
    slide    = _blank(prs)
    _set_bg(slide, TATA_OFF_WHITE)
    insights = d.get("insights", [])
    layout   = _next_layout()

    _header_band(slide, d.get("title", ""), d.get("subtitle", ""))

    ct = HEADER_H + SEP_LINE_H + Inches(0.08)
    ch = SLIDE_H - ct - Inches(0.12)

    if layout == "full":
        img_h  = ch - (Inches(1.9) if insights else Inches(0.1))
        img    = chart_to_png(d, width_in=12.0, height_in=img_h / 914400)
        if img:
            _img(slide, img, Inches(0.5), ct, SLIDE_W - Inches(1.0), img_h)
        if insights:
            _insight_chips(slide, insights,
                           Inches(0.5), ct + img_h + Inches(0.15),
                           SLIDE_W - Inches(1.0), Inches(1.8))

    elif layout == "panel-right":
        cw   = Inches(8.2)
        pw   = SLIDE_W - cw - Inches(0.8)
        pl   = cw + Inches(0.6)
        img  = chart_to_png(d, width_in=8.0, height_in=ch / 914400)
        if img:
            _img(slide, img, Inches(0.4), ct, cw, ch)
        if insights:
            _rect(slide, pl - Inches(0.1), ct, pw + Inches(0.2), ch, TATA_DARK_BLUE)
            _txt(slide, pl, ct + Inches(0.2), pw, Inches(0.4),
                 "KEY INSIGHTS", size=9, bold=True,
                 color=RGBColor(0x03,0xA9,0xF4), font="Calibri")
            for i, ins in enumerate(insights[:4]):
                it = ct + Inches(0.7) + i * Inches(1.1)
                _rect(slide, pl, it, Inches(0.04), Inches(0.85),
                      RGBColor(0x03,0xA9,0xF4))
                _txt(slide, pl + Inches(0.15), it + Inches(0.05),
                     pw - Inches(0.2), Inches(0.9),
                     ins, size=10, color=RGBColor(0xBB,0xDE,0xFB),
                     wrap=True, font="Calibri")

    else:  # dark-sidebar
        sw  = Inches(3.6)
        cl2 = sw + Inches(0.3)
        cw2 = SLIDE_W - cl2 - Inches(0.4)
        _rect(slide, Inches(0), ct, sw, ch, TATA_DARK_BLUE)
        if insights:
            _txt(slide, Inches(0.2), ct + Inches(0.2), sw - Inches(0.3), Inches(0.4),
                 "ANALYSIS", size=9, bold=True,
                 color=RGBColor(0x03,0xA9,0xF4), font="Calibri")
            for i, ins in enumerate(insights[:3]):
                it = ct + Inches(0.75) + i * Inches(1.35)
                _rect(slide, Inches(0.2), it, Inches(0.04), Inches(1.1),
                      RGBColor(0xF5,0x9E,0x0B))
                _txt(slide, Inches(0.36), it + Inches(0.08),
                     sw - Inches(0.5), Inches(1.1),
                     ins, size=10, color=RGBColor(0xBB,0xDE,0xFB),
                     wrap=True, font="Calibri")
        img = chart_to_png(d, width_in=cw2 / 914400, height_in=ch / 914400)
        if img:
            _img(slide, img, cl2, ct, cw2, ch)

    add_bottom_bar(slide, slide_num, total)
    return slide


def build_split(prs: Presentation, d: Dict,
                slide_num: int = None, total: int = None):
    slide    = _blank(prs)
    _set_bg(slide, TATA_OFF_WHITE)
    bullets  = d.get("bullets", [])
    insights = d.get("insights", [])

    _header_band(slide, d.get("title", ""), d.get("subtitle", ""))

    ct = HEADER_H + SEP_LINE_H + Inches(0.08)
    ch = SLIDE_H - ct - Inches(0.12)
    lw = Inches(5.2)

    _rect(slide, Inches(0.4), ct, lw, ch, TATA_DARK_BLUE)
    if bullets:
        _txt(slide, Inches(0.55), ct + Inches(0.15), Inches(0.7), Inches(0.35),
             "HIGHLIGHTS", size=8, bold=True,
             color=RGBColor(0x03,0xA9,0xF4), font="Calibri")
        tb = slide.shapes.add_textbox(Inches(0.65), ct + Inches(0.60),
                                       lw - Inches(0.45), ch - Inches(0.8))
        tf = tb.text_frame; tf.word_wrap = True
        for j, b in enumerate(bullets):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = f"▸  {b}"
            run.font.name  = "Calibri"
            run.font.size  = Pt(11.5)
            run.font.color.rgb = RGBColor(0xBB,0xDE,0xFB)

    rl  = Inches(5.9)
    rw  = SLIDE_W - rl - Inches(0.4)
    img = chart_to_png(d, width_in=rw / 914400, height_in=(ch - Inches(0.3)) / 914400)
    if img:
        _img(slide, img, rl, ct, rw, ch - Inches(0.3))

    if insights:
        _insight_chips(slide, insights,
                       Inches(0.4), ct + ch - Inches(0.1),
                       SLIDE_W - Inches(0.8), Inches(0.8))

    add_bottom_bar(slide, slide_num, total)
    return slide


def build_table(prs: Presentation, d: Dict,
                slide_num: int = None, total: int = None):
    slide    = _blank(prs)
    _set_bg(slide, TATA_OFF_WHITE)
    data     = d.get("data", [])
    insights = d.get("insights", [])

    _header_band(slide, d.get("title", "Data Table"))

    if not data:
        add_bottom_bar(slide, slide_num, total)
        return slide

    if isinstance(data[0], dict):
        cols = list(data[0].keys())
    else:
        cols = ["Value"]
        data = [{"Value": str(r)} for r in data]

    max_rows = min(len(data), 12)
    tt       = Inches(1.2)
    th       = Inches(4.6) if not insights else Inches(3.8)
    tw       = SLIDE_W - Inches(1.0)

    tbl = slide.shapes.add_table(max_rows + 1, len(cols),
                                   Inches(0.5), tt, tw, th).table
    for j, col in enumerate(cols):
        cell = tbl.cell(0, j)
        cell.text = col.replace("_", " ").title()
        cell.fill.solid()
        cell.fill.fore_color.rgb = TATA_DARK_BLUE
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.runs[0] if para.runs else para.add_run()
        run.font.bold = True; run.font.size = Pt(11)
        run.font.color.rgb = TATA_WHITE; run.font.name = "Calibri"

    for i, row in enumerate(data[:max_rows]):
        for j, col in enumerate(cols):
            cell = tbl.cell(i + 1, j)
            cell.text = str(row.get(col, ""))
            cell.fill.solid()
            cell.fill.fore_color.rgb = (TATA_WHITE if i % 2 == 0 else TATA_LIGHT_BLUE)
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            run = para.runs[0] if para.runs else para.add_run()
            run.font.size = Pt(10); run.font.color.rgb = TATA_TEXT_DARK
            run.font.name = "Calibri"

    if insights:
        _insight_chips(slide, insights,
                       Inches(0.5), tt + th + Inches(0.15),
                       SLIDE_W - Inches(1.0), Inches(1.4))

    add_bottom_bar(slide, slide_num, total)
    return slide


def build_timeline(prs: Presentation, d: Dict,
                   slide_num: int = None, total: int = None):
    slide    = _blank(prs)
    _set_bg(slide, TATA_OFF_WHITE)
    data     = d.get("data", [])
    insights = d.get("insights", [])

    _header_band(slide, d.get("title", "Timeline"))

    ct    = HEADER_H + SEP_LINE_H + Inches(0.08)
    img_h = Inches(3.8) if not insights else Inches(3.1)
    img   = timeline_to_png(data, title="", width_in=12.0, height_in=img_h / 914400)
    if img:
        _img(slide, img, Inches(0.4), ct, SLIDE_W - Inches(0.8), img_h)
    if insights:
        _insight_chips(slide, insights,
                       Inches(0.5), ct + img_h + Inches(0.15),
                       SLIDE_W - Inches(1.0), Inches(1.9))

    add_bottom_bar(slide, slide_num, total)
    return slide


def build_bullets(prs: Presentation, d: Dict,
                  slide_num: int = None, total: int = None):
    """
    V12 ── Renders both bullet points AND insight chips on the same slide.
    Especially important for "So What? — Recommended Actions" which
    previously showed empty because insights were not rendered.
    """
    slide    = _blank(prs)
    _set_bg(slide, TATA_OFF_WHITE)
    bullets  = d.get("bullets", d.get("data", []))
    insights = d.get("insights", [])

    _header_band(slide, d.get("title", "Key Insights"))

    has_insights = bool(insights)

    # Shrink the dark panel height when insights are present to make room below
    pt = Inches(1.2)
    ph = (SLIDE_H - pt - Inches(2.1)) if has_insights else (SLIDE_H - pt - Inches(0.4))
    _rect(slide, Inches(0.5), pt, SLIDE_W - Inches(1.0), ph, TATA_DARK_BLUE)

    if bullets:
        items = bullets if (bullets and isinstance(bullets[0], str)) else [
            str(b.get("label", "") or b.get("event", "") or str(b)) for b in bullets
        ]
        tb = slide.shapes.add_textbox(Inches(0.85), pt + Inches(0.35),
                                       SLIDE_W - Inches(1.8), ph - Inches(0.55))
        tf = tb.text_frame; tf.word_wrap = True
        for j, b in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = f"▸   {b}"
            run.font.name  = "Calibri"
            run.font.size  = Pt(13)
            run.font.color.rgb = RGBColor(0xBB,0xDE,0xFB)

    # V12 ── Insight chips below the bullets panel
    if has_insights:
        chips_top = pt + ph + Inches(0.15)
        chips_h   = SLIDE_H - chips_top - Inches(0.18)
        _insight_chips(slide, insights,
                       Inches(0.5), chips_top,
                       SLIDE_W - Inches(1.0), chips_h)

    add_bottom_bar(slide, slide_num, total)
    return slide


def build_thankyou(prs: Presentation, d: Dict):
    """V10 ── Geometric background shapes on Thank You slide."""
    slide = _blank(prs)
    _set_bg(slide, TATA_DARK_BLUE)

    for (l, t, w, h, rot, color) in [
        (Inches(-1.5), Inches(-1.0), Inches(5.0), Inches(4.0),  25, RGBColor(0x1A,0x5C,0xC0)),
        (Inches(10.5), Inches(5.5),  Inches(5.0), Inches(4.0), -20, RGBColor(0x0A,0x3A,0x80)),
        (Inches(11.0), Inches(-0.5), Inches(3.5), Inches(3.5),  15, RGBColor(0x03,0x7A,0xB5)),
    ]:
        s = slide.shapes.add_shape(1, l, t, w, h)
        s.fill.solid(); s.fill.fore_color.rgb = color
        s.line.fill.background(); s.rotation = rot

    _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), TATA_BLUE)

    _txt(slide, Inches(1.5), Inches(2.1), Inches(10.3), Inches(1.8),
         "Thank You", size=56, bold=True, color=TATA_WHITE,
         align=PP_ALIGN.CENTER, font="Calibri")
    _rect(slide, Inches(4.2), Inches(3.95), Inches(4.9), Inches(0.06),
          RGBColor(0x03,0xA9,0xF4))

    sub = d.get("subtitle", "")
    if sub:
        _txt(slide, Inches(2.0), Inches(4.2), Inches(9.3), Inches(1.0),
             sub, size=16, color=RGBColor(0xBB,0xDE,0xFB),
             align=PP_ALIGN.CENTER, font="Calibri")

    add_logo(slide, size_inches=1.0)
    return slide

# ─── DISPATCH TABLE ───────────────────────────────────────────────────────────

_PAGINATED = {"kpi", "chart", "split", "table", "timeline", "bullets"}

_BUILDERS = {
    "cover":    build_cover,
    "section":  build_section,
    "kpi":      build_kpi,
    "chart":    build_chart,
    "split":    build_split,
    "table":    build_table,
    "timeline": build_timeline,
    "bullets":  build_bullets,
    "thankyou": build_thankyou,
}

# ─── S1: SANITIZER ────────────────────────────────────────────────────────────

def _sanitize(obj, depth: int = 0):
    """S1 ── Recursively truncate strings; cap depth to prevent DoS."""
    if depth > 5:
        return {}
    if isinstance(obj, dict):
        return {k: _sanitize(v, depth + 1)
                for k, v in obj.items()
                if isinstance(k, str) and len(k) <= 100}
    if isinstance(obj, list):
        return [_sanitize(i, depth + 1) for i in obj[:200]]
    if isinstance(obj, str):
        return obj[:500]
    return obj

# ─── MAIN ENTRY ───────────────────────────────────────────────────────────────

def generate_pptx(slides: List[Dict], watermark: str = None) -> bytes:
    """
    Build PPTX bytes from slide dicts.

    Args:
        slides:    Output from ppt_brain.generate_slide_plan().
        watermark: Optional 'DRAFT' or 'CONFIDENTIAL' stamp on every slide (E4).
    """
    _reset_layouts()

    slides = [_sanitize(s) for s in slides]

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    if not any(s.get("slide_type") == "cover" for s in slides):
        slides = [{"slide_type": "cover", "title": "Presentation"}] + slides
    if not any(s.get("slide_type") == "thankyou" for s in slides):
        slides = slides + [{"slide_type": "thankyou",
                             "subtitle": "Tata Chemicals Limited"}]

    renderable    = [s for s in slides if has_enough_data(s)]
    total         = len(renderable)
    rendered      = 0
    section_count = 0

    for sd in slides:
        st = sd.get("slide_type", "bullets").lower()

        # NO_FUNNEL ── skip any pyramid/funnel slides that may have leaked through
        if st == "pyramid":
            logger.info(f"Skipping pyramid slide '{sd.get('title','')}' — funnel removed")
            continue

        if not has_enough_data(sd):
            logger.warning(f"Skipping '{sd.get('title','')}' — insufficient data")
            continue

        rendered += 1
        try:
            if st == "cover":
                built = build_cover(prs, sd)
            elif st == "section":
                section_count += 1
                built = build_section(prs, sd, section_count)
            elif st == "thankyou":
                built = build_thankyou(prs, sd)
            elif st in _PAGINATED:
                built = _BUILDERS[st](prs, sd, slide_num=rendered, total=total)
            else:
                logger.warning(f"Unknown slide_type '{st}' → bullets fallback")
                built = build_bullets(prs, sd, rendered, total)

            if watermark and built is not None:
                _watermark(built, watermark)

        except Exception as e:
            logger.error(f"Build error slide '{st}': {e}", exc_info=True)
            try:
                build_bullets(prs, {
                    "title":   sd.get("title", "Slide"),
                    "bullets": [str(x) for x in sd.get("data", [])],
                }, rendered, total)
            except Exception as fe:
                logger.error(f"Fallback also failed: {fe}")

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()