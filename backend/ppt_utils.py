"""
ppt_utils.py  ── v2.1
Helpers for PPT generation.

UPGRADES:
  V3  Transparent chart backgrounds
  V9  Slide number + progress dots support
  C1  Waterfall chart
  C2  Bullet (actual vs target) chart
  C3  Treemap chart
  C4  Scatter / quadrant chart
  C5  Dual-axis combo chart
  C6  Semantic color coding (Closed Won=green, Closed Lost=red, etc.)
  C7  Value labels on every bar/line data point
  C8  Chart title removed from matplotlib (slide title covers it)
  A2  SOQL result cache (5-min TTL, keyed by query hash)
  A4  Python-side anomaly detection (z-scores, MoM deltas, Pareto)
  FIX  Zero-sum guard: skip charts with no positive data instead of crashing
"""

import io
import os
import hashlib
import logging
import time
from typing import List, Dict, Any, Optional, Tuple

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

logger = logging.getLogger(__name__)

# ─── TEMPLATE PATH ────────────────────────────────────────────────────────────
TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), "Tata_Chemicals_PPT_Template_FY_25.pptx")

# ─── BRAND PALETTE ────────────────────────────────────────────────────────────
TATA_BLUE       = RGBColor(0x15, 0x65, 0xC0)
TATA_DARK_BLUE  = RGBColor(0x0D, 0x47, 0xA1)
TATA_LIGHT_BLUE = RGBColor(0xE3, 0xF0, 0xFF)
TATA_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
TATA_OFF_WHITE  = RGBColor(0xF5, 0xF7, 0xFA)
TATA_TEXT_DARK  = RGBColor(0x1E, 0x29, 0x3B)
TATA_TEXT_MID   = RGBColor(0x44, 0x57, 0x6D)
TATA_TEXT_MUTED = RGBColor(0x64, 0x74, 0x8B)
TATA_ACCENT1    = RGBColor(0x03, 0xA9, 0xF4)
TATA_SUCCESS    = RGBColor(0x43, 0xA0, 0x47)
TATA_WARNING    = RGBColor(0xFB, 0x8C, 0x00)
TATA_DANGER     = RGBColor(0xE5, 0x39, 0x35)
TATA_TEAL       = RGBColor(0x00, 0x83, 0x8F)

# Matplotlib hex palette
MPL_PALETTE = [
    "#1565C0", "#03A9F4", "#0D47A1", "#43A047",
    "#FB8C00", "#E53935", "#7B1FA2", "#00838F",
    "#F57F17", "#37474F",
]

# C6 ── Semantic stage/status colors
STAGE_COLORS: Dict[str, str] = {
    "closed won":    "#43A047",
    "won":           "#43A047",
    "closed lost":   "#E53935",
    "lost":          "#E53935",
    "commercial":    "#1565C0",
    "qualification": "#03A9F4",
    "sampling":      "#FB8C00",
    "negotiation":   "#7B1FA2",
    "proposal":      "#00838F",
    "prospecting":   "#F57F17",
}

def semantic_color(label: str, fallback: str) -> str:
    return STAGE_COLORS.get(str(label).lower().strip(), fallback)

# ─── SLIDE DIMENSIONS ─────────────────────────────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)

# ─── A2: SOQL RESULT CACHE ────────────────────────────────────────────────────
_soql_cache: Dict[str, Tuple[list, float]] = {}
CACHE_TTL = 300

def cache_get(query: str) -> Optional[list]:
    key = hashlib.md5(query.strip().lower().encode()).hexdigest()
    entry = _soql_cache.get(key)
    if entry and (time.time() - entry[1]) < CACHE_TTL:
        logger.info(f"SOQL cache HIT [{key[:8]}]")
        return entry[0]
    return None

def cache_set(query: str, results: list) -> None:
    key = hashlib.md5(query.strip().lower().encode()).hexdigest()
    _soql_cache[key] = (results, time.time())

def cache_clear_expired() -> None:
    now = time.time()
    stale = [k for k, (_, ts) in list(_soql_cache.items()) if now - ts > CACHE_TTL]
    for k in stale:
        del _soql_cache[k]

# ─── LOGO ─────────────────────────────────────────────────────────────────────
_logo_bytes: Optional[bytes] = None

def get_logo_bytes() -> Optional[bytes]:
    global _logo_bytes
    if _logo_bytes is not None:
        return _logo_bytes
    try:
        prs = Presentation(TEMPLATE_PATH)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == 13:
                    _logo_bytes = shape.image.blob
                    logger.info("Logo extracted from template.")
                    return _logo_bytes
    except Exception as e:
        logger.warning(f"Could not extract logo: {e}")
    return None

def add_logo(slide, size_inches: float = 0.9) -> None:
    logo = get_logo_bytes()
    if not logo:
        return
    try:
        w = Inches(size_inches * 2.0)
        h = Inches(size_inches * 0.55)
        slide.shapes.add_picture(io.BytesIO(logo),
                                  SLIDE_W - w - Inches(0.25),
                                  SLIDE_H - h - Inches(0.18), w, h)
    except Exception as e:
        logger.warning(f"add_logo failed: {e}")

def add_logo_top_right(slide, size_inches: float = 1.0) -> None:
    logo = get_logo_bytes()
    if not logo:
        return
    try:
        w = Inches(size_inches * 2.0)
        h = Inches(size_inches * 0.55)
        slide.shapes.add_picture(io.BytesIO(logo),
                                  SLIDE_W - w - Inches(0.3), Inches(0.22), w, h)
    except Exception as e:
        logger.warning(f"add_logo_top_right failed: {e}")

def add_bottom_bar(slide, slide_num: int = None, total_slides: int = None) -> None:
    """V9 ── Red bottom bar matching Tata Chemicals template with slide number + progress dots."""
    from pptx.enum.text import PP_ALIGN
    TATA_RED_BAR = RGBColor(0xE0, 0x20, 0x3C)

    bar_h = Inches(0.057)
    bar_t = SLIDE_H - bar_h
    s = slide.shapes.add_shape(1, Inches(0), bar_t, SLIDE_W, bar_h)
    s.fill.solid()
    s.fill.fore_color.rgb = TATA_RED_BAR
    s.line.fill.background()

    if slide_num is not None:
        tb = slide.shapes.add_textbox(SLIDE_W - Inches(0.55),
                                       bar_t - Inches(0.24),
                                       Inches(0.45), Inches(0.22))
        tf = tb.text_frame
        p  = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(slide_num)
        run.font.size = Pt(7)
        run.font.color.rgb = RGBColor(0x44, 0x57, 0x6D)
        run.font.name = "Calibri"

    if slide_num is not None and total_slides and total_slides > 1:
        dot_sz  = Inches(0.07)
        dot_gap = Inches(0.12)
        total_w = total_slides * dot_sz + (total_slides - 1) * dot_gap
        start_x = (SLIDE_W - total_w) / 2
        dot_y   = bar_t - Inches(0.14)
        for i in range(total_slides):
            dl = start_x + i * (dot_sz + dot_gap)
            d  = slide.shapes.add_shape(9, dl, dot_y, dot_sz, dot_sz)
            d.fill.solid()
            d.fill.fore_color.rgb = (TATA_DARK_BLUE if i == slide_num - 1
                                     else RGBColor(0xCC, 0xCC, 0xCC))
            d.line.fill.background()

# ─── CHART STYLE HELPER ───────────────────────────────────────────────────────

def _apply_clean_style(ax, fig) -> None:
    """V3 ── Transparent background so charts blend into any slide colour."""
    fig.patch.set_alpha(0)
    ax.set_facecolor("none")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#CBD5E1")
    ax.spines["bottom"].set_color("#CBD5E1")
    ax.tick_params(colors="#44576D", labelsize=9)
    ax.yaxis.label.set_color("#44576D")
    ax.xaxis.label.set_color("#44576D")
    ax.grid(axis="y", color="#CBD5E1", linewidth=0.7, linestyle="--", alpha=0.6)
    ax.set_axisbelow(True)
    plt.rcParams["font.family"] = "DejaVu Sans"

def _value_labels(ax, bars, values) -> None:
    """C7 ── Value labels above every bar. Guards against zero-sum."""
    max_v = max((abs(v) for v in values), default=1) or 1  # FIX: never 0
    for bar, val in zip(bars, values):
        ax.text(
            bar.get_x() + bar.get_width() / 2,
            bar.get_height() + max_v * 0.012,
            f"{val:,.0f}",
            ha="center", va="bottom",
            fontsize=8, color="#1E293B", fontweight="bold"
        )

def _set_xlabels(ax, labels) -> None:
    ax.set_xticks(range(len(labels)))
    ax.set_xticklabels(labels,
                       rotation=20 if len(labels) > 5 else 0,
                       ha="right", fontsize=9)

# ─── ZERO-SUM GUARD ───────────────────────────────────────────────────────────

def _has_plottable_data(values: list) -> bool:
    """
    FIX ── Return False when all values are zero or negative for chart types
    that would cause matplotlib division-by-zero or invisible renders.
    """
    return any(v != 0 for v in values)

# ─── INDIVIDUAL CHART RENDERERS ───────────────────────────────────────────────

def _chart_bar(ax, fig, labels, values, colors) -> None:
    bars = ax.bar(labels, values, color=colors, width=0.55, zorder=3, linewidth=0)
    _value_labels(ax, bars, values)
    _set_xlabels(ax, labels)

def _chart_line(ax, fig, labels, values) -> None:
    if any(" - Target" in l or " - Actual" in l for l in labels):
        t_pts, a_pts = {}, {}
        for lbl, val in zip(labels, values):
            if " - Target" in lbl:
                t_pts[lbl.replace(" - Target", "")] = val
            elif " - Actual" in lbl:
                a_pts[lbl.replace(" - Actual", "")] = val
        dims = sorted(set(list(t_pts) + list(a_pts)))
        ax.plot(dims, [t_pts.get(d, 0) for d in dims],
                marker="o", lw=2.5, color="#1565C0", label="Target", ms=7)
        ax.plot(dims, [a_pts.get(d, 0) for d in dims],
                marker="s", lw=2.5, color="#43A047", label="Actual", ms=7, ls="--")
        ax.legend(fontsize=9, framealpha=0, labelcolor="#44576D")
        _set_xlabels(ax, dims)
    else:
        ax.plot(labels, values, marker="o", lw=2.5, color="#1565C0", ms=7, zorder=3)
        ax.fill_between(range(len(labels)), values, alpha=0.10, color="#1565C0")
        max_v = max(values) if values else 1
        max_v = max_v or 1  # FIX: guard against zero max
        for i, (x, y) in enumerate(zip(range(len(labels)), values)):
            ax.text(x, y + max_v * 0.015, f"{y:,.0f}",
                    ha="center", va="bottom", fontsize=7.5,
                    color="#1E293B", fontweight="bold")
        _set_xlabels(ax, labels)

def _chart_pie(ax, fig, labels, values, donut: bool = False) -> None:
    PIE = ["#1565C0","#03A9F4","#00838F","#43A047",
           "#FB8C00","#E53935","#7B1FA2","#F57F17",
           "#0288D1","#2E7D32","#AD1457","#37474F"]
    wc    = (PIE * ((len(labels) // len(PIE)) + 1))[:len(labels)]
    total = sum(values) if sum(values) > 0 else 1  # FIX: guard zero-sum
    leg   = [f"{l}  —  {v:,.1f}  ({v/total*100:.1f}%)" for l, v in zip(labels, values)]
    wedges, _ = ax.pie(values, colors=wc, startangle=90, counterclock=False,
                       wedgeprops={"linewidth": 2.5, "edgecolor": "#FFFFFF"})
    if donut:
        ax.add_patch(plt.Circle((0, 0), 0.65, fc="none"))
        ax.text(0, 0, f"{total:,.0f}", ha="center", va="center",
                fontsize=14, fontweight="bold", color="#1E293B")
    ax.legend(wedges, leg, loc="center left", bbox_to_anchor=(1.02, 0.5),
              fontsize=8.5, frameon=False)

def _chart_waterfall(ax, fig, labels, values) -> None:
    """C1 ── Waterfall: cumulative positive/negative bars."""
    running = 0
    bottoms, bar_vals, colors = [], [], []
    for i, val in enumerate(values):
        if i == 0 or i == len(values) - 1:
            bottoms.append(0); bar_vals.append(val); colors.append("#1565C0")
        else:
            bottoms.append(running); bar_vals.append(val)
            colors.append("#43A047" if val >= 0 else "#E53935")
        if 0 < i < len(values) - 1:
            running += val
    bars = ax.bar(labels, bar_vals, bottom=bottoms, color=colors, width=0.55, zorder=3, linewidth=0)
    max_abs = max((abs(v) for v in values), default=1) or 1  # FIX
    for bar, bv, bt in zip(bars, bar_vals, bottoms):
        ax.text(bar.get_x() + bar.get_width() / 2,
                bt + bv + max_abs * 0.012,
                f"{bv:+,.0f}", ha="center", va="bottom",
                fontsize=8, color="#1E293B", fontweight="bold")
    ax.axhline(0, color="#CBD5E1", linewidth=0.8)
    _set_xlabels(ax, labels)
    ax.legend(handles=[
        mpatches.Patch(color="#1565C0", label="Base/Total"),
        mpatches.Patch(color="#43A047", label="Increase"),
        mpatches.Patch(color="#E53935", label="Decrease"),
    ], fontsize=8, framealpha=0, labelcolor="#44576D", loc="upper right")

def _chart_bullet(ax, fig, labels, data) -> None:
    """C2 ── Bullet chart: actual vs target horizontal bars."""
    pairs: Dict[str, Dict] = {}
    for item in data:
        lbl = str(item.get("label", ""))
        val = float(item.get("value", 0))
        if " - Actual" in lbl:
            pairs.setdefault(lbl.replace(" - Actual", ""), {})["actual"] = val
        elif " - Target" in lbl:
            pairs.setdefault(lbl.replace(" - Target", ""), {})["target"] = val
        else:
            pairs.setdefault(lbl, {})["actual"] = val
    keys    = list(pairs.keys())
    actuals = [pairs[k].get("actual", 0) for k in keys]
    targets = [pairs[k].get("target", 0) for k in keys]
    y_pos   = range(len(keys))
    ax.barh(list(y_pos), targets, height=0.55, color="#CBD5E1", zorder=2, label="Target")
    ax.barh(list(y_pos), actuals, height=0.30, color="#1565C0", zorder=3, label="Actual")
    max_target = max(targets + [1]) or 1  # FIX: guard zero-sum
    for i, (a, t) in enumerate(zip(actuals, targets)):
        pct   = (a / t * 100) if t else 0
        color = "#43A047" if pct >= 100 else ("#FB8C00" if pct >= 75 else "#E53935")
        ax.text(max(a, t) + max_target * 0.02, i,
                f"{pct:.0f}%", va="center", fontsize=8, color=color, fontweight="bold")
    ax.set_yticks(list(y_pos))
    ax.set_yticklabels(keys, fontsize=9)
    ax.legend(fontsize=8, framealpha=0, labelcolor="#44576D", loc="lower right")
    ax.spines["left"].set_visible(False)
    ax.tick_params(left=False)

def _chart_treemap(ax, fig, labels, values) -> None:
    """C3 ── Treemap (squarify if available, bar fallback otherwise)."""
    try:
        import squarify
        ax.set_xlim(0, 100); ax.set_ylim(0, 100); ax.axis("off")
        total   = sum(values) if sum(values) > 0 else 1  # FIX
        colors  = (MPL_PALETTE * ((len(labels) // len(MPL_PALETTE)) + 1))[:len(labels)]
        rects   = squarify.squarify(values, x=0, y=0, dx=100, dy=100)
        for rect, label, val, color in zip(rects, labels, values, colors):
            ax.add_patch(plt.Rectangle((rect["x"], rect["y"]),
                                        rect["dx"], rect["dy"],
                                        color=color, ec="white", lw=2))
            if rect["dx"] > 8 and rect["dy"] > 5:
                ax.text(rect["x"] + rect["dx"] / 2, rect["y"] + rect["dy"] / 2,
                        f"{label}\n{val:,.0f}\n({val/total*100:.1f}%)",
                        ha="center", va="center",
                        fontsize=max(7, min(11, rect["dx"] * 0.5)),
                        color="white", fontweight="bold")
    except ImportError:
        colors = (MPL_PALETTE * ((len(labels) // len(MPL_PALETTE)) + 1))[:len(labels)]
        bars   = ax.bar(labels, values, color=colors, width=0.55)
        _value_labels(ax, bars, values)
        _set_xlabels(ax, labels)

def _chart_scatter(ax, fig, data) -> None:
    """C4 ── Scatter / quadrant chart (amount vs probability)."""
    xs, ys, lbls = [], [], []
    for item in data:
        xs.append(float(item.get("x", item.get("amount", item.get("value", 0)))))
        ys.append(float(item.get("y", item.get("probability", 50))))
        lbls.append(str(item.get("label", "")))
    mid_x = float(np.median(xs)) if xs else 0
    mid_y = float(np.median(ys)) if ys else 50
    ax.axvline(mid_x, color="#CBD5E1", lw=1, ls="--", alpha=0.7)
    ax.axhline(mid_y, color="#CBD5E1", lw=1, ls="--", alpha=0.7)
    if xs and ys:
        max_x = max(xs) or 1  # FIX
        max_y = max(ys) or 1  # FIX
        ax.fill_betweenx([mid_y, max_y * 1.1], mid_x, max_x * 1.1,
                         alpha=0.04, color="#43A047")
        ax.fill_betweenx([0, mid_y], 0, mid_x, alpha=0.04, color="#E53935")
    colors = (MPL_PALETTE * ((len(xs) // len(MPL_PALETTE)) + 1))[:len(xs)]
    ax.scatter(xs, ys, c=colors, s=80, zorder=3, alpha=0.85)
    for x, y, lbl in zip(xs, ys, lbls):
        ax.annotate(lbl, (x, y), textcoords="offset points", xytext=(6, 4),
                    fontsize=7.5, color="#1E293B")
    if xs and ys:
        xm = (max(xs) or 1) * 1.05
        ym = (max(ys) or 1) * 1.05
        ax.text(xm * 0.98, ym * 0.97, "High Value\nHigh Prob",
                ha="right", va="top", fontsize=7, color="#43A047", alpha=0.7)
        ax.text(xm * 0.02, ym * 0.97, "Low Value\nHigh Prob",
                ha="left",  va="top", fontsize=7, color="#FB8C00", alpha=0.7)

def _chart_combo(ax, fig, labels, values, data) -> None:
    """C5 ── Dual-axis combo: bars (primary) + line (secondary value2)."""
    values2 = [float(d.get("value2", 0)) for d in data]
    colors  = [semantic_color(l, MPL_PALETTE[i % len(MPL_PALETTE)])
               for i, l in enumerate(labels)]
    bars = ax.bar(labels, values, color=colors, width=0.55, zorder=3, alpha=0.85, linewidth=0)
    _set_xlabels(ax, labels)
    ax2 = ax.twinx()
    ax2.plot(range(len(labels)), values2, marker="o", lw=2.5,
             color="#E53935", ms=7, zorder=4)
    ax2.tick_params(colors="#E53935", labelsize=9)
    ax2.spines["top"].set_visible(False)
    ax2.spines["right"].set_color("#E53935")
    ax2.set_facecolor("none")
    ax.legend(handles=[
        mpatches.Patch(color="#1565C0", label="Primary (bars)"),
        mpatches.Patch(color="#E53935", label="Secondary (line)"),
    ], fontsize=8, framealpha=0, labelcolor="#44576D", loc="upper left")

# ─── MAIN CHART DISPATCHER ────────────────────────────────────────────────────

def chart_to_png(chart_config: Dict[str, Any],
                 width_in: float = 8.5,
                 height_in: float = 4.2) -> Optional[bytes]:
    """
    Convert chart config → transparent PNG bytes.
    Supports: bar, line, pie, donut, waterfall, bullet, treemap, scatter, combo.
    FIX: Returns None (instead of crashing) when all data values are zero.
    """
    try:
        chart_type = chart_config.get("chart_type", "bar").lower()
        x_label    = chart_config.get("x_axis", "")
        y_label    = chart_config.get("y_axis", "")
        data       = chart_config.get("data", [])
        if not data:
            return None

        labels = [str(d.get("label", "")) for d in data]
        values = [float(d.get("value", 0)) for d in data]

        # FIX ── Zero-sum guard: skip rendering for chart types that
        # would cause division-by-zero or invisible output.
        if chart_type in ("pie", "donut", "treemap") and not _has_plottable_data(values):
            logger.warning(f"chart_to_png: all-zero values for '{chart_type}' — returning None")
            return None

        colors = [semantic_color(l, MPL_PALETTE[i % len(MPL_PALETTE)])
                  for i, l in enumerate(labels)]

        wide = chart_type in ("pie", "donut", "treemap")
        fig, ax = plt.subplots(figsize=(width_in + (2.5 if wide else 0), height_in))
        _apply_clean_style(ax, fig)

        if chart_type == "bar":
            _chart_bar(ax, fig, labels, values, colors)
        elif chart_type == "line":
            _chart_line(ax, fig, labels, values)
        elif chart_type == "pie":
            _chart_pie(ax, fig, labels, values, donut=False)
        elif chart_type == "donut":
            _chart_pie(ax, fig, labels, values, donut=True)
        elif chart_type == "waterfall":
            _chart_waterfall(ax, fig, labels, values)
        elif chart_type == "bullet":
            _chart_bullet(ax, fig, labels, data)
        elif chart_type == "treemap":
            _chart_treemap(ax, fig, labels, values)
        elif chart_type == "scatter":
            _chart_scatter(ax, fig, data)
        elif chart_type == "combo":
            _chart_combo(ax, fig, labels, values, data)
        else:
            _chart_bar(ax, fig, labels, values, colors)

        if x_label and chart_type not in ("pie", "donut", "treemap", "scatter"):
            ax.set_xlabel(x_label, fontsize=9)
        if y_label and chart_type not in ("pie", "donut", "treemap", "scatter"):
            ax.set_ylabel(y_label, fontsize=9)

        if wide:
            plt.tight_layout(pad=0.5, rect=[0, 0, 0.72, 1])
        else:
            plt.tight_layout(pad=0.8)

        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=180, bbox_inches="tight", transparent=True)
        plt.close(fig)
        buf.seek(0)
        return buf.read()

    except Exception as e:
        logger.error(f"chart_to_png error: {e}")
        plt.close("all")
        return None

# ─── PYRAMID ──────────────────────────────────────────────────────────────────

def pyramid_to_png(data: List[Dict], title: str = "",
                   width_in: float = 9.0,
                   height_in: float = 5.0) -> Optional[bytes]:
    """
    FIX_PYRAMID ── Log-normalized bar widths when value spread exceeds 50x.
    Conversion % only shown when next stage ≤ current (true funnel drop).
    FIX ── Guards against division by zero when all funnel values are 0.
    """
    try:
        if not data:
            return None
        labels = [str(d.get("label", "")) for d in data]
        values = [float(d.get("value", 0)) for d in data]
        n      = len(labels)
        fig, ax = plt.subplots(figsize=(width_in, height_in))
        fig.patch.set_facecolor("#0F172A")
        ax.set_facecolor("#0F172A")
        ax.axis("off")
        ax.set_xlim(0, 10)
        ax.set_ylim(0, n + 0.5)
        palette = ["#1E40AF","#1D4ED8","#2563EB","#3B82F6",
                   "#60A5FA","#93C5FD","#BFDBFE","#DBEAFE"]
        palette = (palette * ((n // len(palette)) + 1))[:n]

        max_val  = max(values) if max(values) > 0 else 1  # FIX
        pos_vals = [v for v in values if v > 0]
        min_val  = min(pos_vals) if pos_vals else 1
        use_log  = (max_val / min_val) > 50 if min_val > 0 else False

        def _norm_width(v: float) -> float:
            if use_log:
                import math
                log_v   = math.log1p(max(v, 0))
                log_max = math.log1p(max_val)
                return 2.0 + 6.5 * (log_v / log_max) if log_max > 0 else 2.0
            return 2.0 + 6.5 * (v / max_val)  # max_val guaranteed > 0 above

        for i, (label, val) in enumerate(zip(labels, values)):
            y_top = n - i
            y_bot = y_top - 1.0 + 0.12
            wf    = _norm_width(val)
            xl, xr = (10 - wf) / 2, (10 + wf) / 2
            if i < n - 1:
                nw  = _norm_width(values[i + 1])
                xlb = (10 - nw) / 2; xrb = (10 + nw) / 2
            else:
                xlb = xl + 0.3; xrb = xr - 0.3
            ax.add_patch(plt.Polygon(
                [[xl, y_top],[xr, y_top],[xrb, y_bot],[xlb, y_bot]],
                closed=True, facecolor=palette[i],
                edgecolor="#0F172A", linewidth=1.5, zorder=2))
            ax.text(5, (y_top + y_bot) / 2,
                    f"{label}  ●  {val:,.0f}",
                    ha="center", va="center",
                    fontsize=11, fontweight="bold", color="white", zorder=3)
            if i < n - 1 and values[i] > 0 and values[i + 1] <= values[i]:
                conv = values[i + 1] / values[i] * 100
                ax.annotate(f"{conv:.0f}%",
                            xy=(xrb + 0.1, y_bot - 0.05),
                            xytext=(xr + 0.15, y_top - 0.5),
                            fontsize=9, color="#F59E0B", fontweight="bold",
                            ha="left", va="center", zorder=4,
                            arrowprops=dict(arrowstyle="->", color="#F59E0B", lw=1.2))
        if title:
            ax.text(5, n + 0.35, title, ha="center", va="center",
                    fontsize=13, fontweight="bold", color="white", zorder=4)
        plt.tight_layout(pad=0.3)
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=160, bbox_inches="tight",
                    facecolor=fig.get_facecolor())
        plt.close(fig)
        buf.seek(0)
        return buf.read()
    except Exception as e:
        logger.error(f"pyramid_to_png error: {e}")
        plt.close("all")
        return None

# ─── TIMELINE ─────────────────────────────────────────────────────────────────

def timeline_to_png(data: List[Dict], title: str = "",
                    width_in: float = 11.5,
                    height_in: float = 3.8) -> Optional[bytes]:
    try:
        if not data:
            return None
        n = len(data)
        fig, ax = plt.subplots(figsize=(width_in, height_in))
        fig.patch.set_facecolor("#0F172A")
        ax.set_facecolor("#0F172A")
        ax.axis("off")
        ax.set_xlim(0, 1)
        ax.set_ylim(0, 1)
        xs     = np.linspace(0.07, 0.93, n)
        y_line = 0.50
        ax.plot([0.05, 0.95], [y_line, y_line],
                color="#3B82F6", linewidth=3.5, zorder=1, solid_capstyle="round")
        ax.plot([0.05, 0.95], [y_line, y_line],
                color="#60A5FA", linewidth=8, alpha=0.15, zorder=0)
        for i, (x, item) in enumerate(zip(xs, data)):
            label = str(item.get("label", ""))
            event = str(item.get("event", ""))
            above = (i % 2 == 0)
            y_text = y_line + 0.32 if above else y_line - 0.32
            y_ev   = y_line - 0.26 if above else y_line + 0.26
            y_se   = y_line + 0.09 if above else y_line - 0.09
            y_ss   = y_text - (0.04 if above else -0.04)
            ax.plot([x, x], [y_se, y_ss], color="#475569", lw=1.2, zorder=2, ls="--")
            ax.plot(x, y_line, "o", color="#F59E0B", ms=22, zorder=3, alpha=0.25)
            ax.plot(x, y_line, "o", color="#FEF3C7", ms=18, zorder=4)
            ax.plot(x, y_line, "o", color="#F59E0B", ms=14, zorder=5)
            ax.text(x, y_line, str(i + 1), ha="center", va="center",
                    fontsize=8, fontweight="bold", color="#1E293B", zorder=6)
            ax.text(x, y_text, label, ha="center",
                    va="bottom" if above else "top",
                    fontsize=9, fontweight="bold", color="#FCD34D", zorder=6)
            ax.text(x, y_ev, event, ha="center",
                    va="top" if above else "bottom",
                    fontsize=7.5, color="#93C5FD", zorder=6, multialignment="center")
        plt.tight_layout(pad=0.4)
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=160, bbox_inches="tight",
                    facecolor=fig.get_facecolor())
        plt.close(fig)
        buf.seek(0)
        return buf.read()
    except Exception as e:
        logger.error(f"timeline_to_png error: {e}")
        plt.close("all")
        return None

# ─── A4: ANOMALY DETECTION ────────────────────────────────────────────────────

def detect_anomalies(data: List[Dict], value_key: str = "value") -> Dict:
    """
    A4 ── Statistical pre-analysis before LLM insight generation.
    Flags: z-score anomalies, MoM deltas, Pareto concentration, peak/trough.
    """
    if not data or len(data) < 3:
        return {}
    try:
        values = [float(d.get(value_key, 0)) for d in data]
        labels = [str(d.get("label", i)) for i, d in enumerate(data)]
        mean   = float(np.mean(values))
        std    = float(np.std(values))
        if std == 0:
            return {"all_equal": True, "mean": mean}
        z_scores  = [(v - mean) / std for v in values]
        anomalies = [
            {"label": l, "value": v, "z_score": round(z, 2)}
            for l, v, z in zip(labels, values, z_scores) if abs(z) > 2.0
        ]
        mom_deltas = []
        for i in range(1, len(values)):
            if values[i - 1] != 0:
                delta = (values[i] - values[i - 1]) / abs(values[i - 1]) * 100
                mom_deltas.append({"from": labels[i - 1], "to": labels[i],
                                   "delta_pct": round(delta, 1)})
        sorted_v    = sorted(values, reverse=True)
        cumulative  = list(np.cumsum(sorted_v))
        total_sum   = cumulative[-1] if cumulative[-1] != 0 else 1
        pareto_idx  = next((i for i, c in enumerate(cumulative) if c / total_sum >= 0.8), None)
        pareto_count = pareto_idx + 1 if pareto_idx is not None else None
        pareto_pct   = round(pareto_count / len(values) * 100, 0) if pareto_count else None
        return {
            "mean":         round(mean, 2),
            "std":          round(std, 2),
            "max_label":    labels[values.index(max(values))],
            "max_value":    max(values),
            "min_label":    labels[values.index(min(values))],
            "min_value":    min(values),
            "anomalies":    anomalies,
            "mom_deltas":   mom_deltas,
            "pareto_count": pareto_count,
            "pareto_pct":   pareto_pct,
            "total":        round(sum(values), 2),
        }
    except Exception as e:
        logger.warning(f"detect_anomalies failed: {e}")
        return {}

# ─── DATA QUALITY ─────────────────────────────────────────────────────────────

def has_enough_data(slide_json: Dict) -> bool:
    slide_type = slide_json.get("slide_type", "")
    data       = slide_json.get("data", [])
    if slide_type in ("cover", "section", "thankyou", "bullets"):
        return True
    if slide_type == "kpi":
        return len(data) >= 1
    if slide_type in ("chart", "split"):
        return len(data) >= 3
    if slide_type == "timeline":
        return len(data) >= 2
    if slide_type == "pyramid":
        return len(data) >= 2
    if slide_type == "table":
        return len(data) >= 1
    return len(data) >= 1